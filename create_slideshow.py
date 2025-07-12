import os
import shutil
from PIL import Image, ExifTags
from pptx import Presentation
from pptx.util import Inches

def create_slideshow_from_folder():
    """
    Automates the creation of a PowerPoint slideshow from a folder of images.
    This script verifies, rotates, resizes, and standardizes images before
    adding them to a widescreen presentation.
    """
    # --- 1. Get User Input and Set Up Paths ---
    source_folder = input("Enter the full path to the folder containing your images: ")

    if not os.path.isdir(source_folder):
        print(f"\n❌ Error: The folder '{source_folder}' was not found.")
        return

    # Create a temporary folder for processed images inside the source folder.
    processed_folder = os.path.join(source_folder, "processed_images")
    if os.path.exists(processed_folder):
        shutil.rmtree(processed_folder) # Clean up old runs
    os.makedirs(processed_folder)

    # Define the output file path on the user's Desktop.
    output_filename = "Generated_Slideshow.pptx"
    output_path = os.path.join(os.path.expanduser("~"), "Desktop", output_filename)
    
    print(f"\nProcessing images from: {source_folder}")
    print(f"Temporary files will be stored in: {processed_folder}")

    try:
        # --- 2. Process All Images ---
        valid_image_extensions = ('.png', '.jpg', '.jpeg', '.gif', '.bmp')
        image_files = sorted([f for f in os.listdir(source_folder) if f.lower().endswith(valid_image_extensions)])

        if not image_files:
            print("\n❌ No valid image files found in the specified folder.")
            return

        print(f"\nFound {len(image_files)} images. Starting processing...")

        for filename in image_files:
            try:
                filepath = os.path.join(source_folder, filename)
                with Image.open(filepath) as img:
                    # --- Fix Orientation (from EXIF data) ---
                    exif = img.getexif()
                    orientation_key = None
                    # Find the orientation tag key
                    for key, val in ExifTags.TAGS.items():
                        if val == 'Orientation':
                            orientation_key = key
                            break
                    
                    if orientation_key and orientation_key in exif:
                        orientation = exif[orientation_key]
                        if orientation == 3:
                            img = img.rotate(180, expand=True)
                        elif orientation == 6:
                            img = img.rotate(270, expand=True)
                        elif orientation == 8:
                            img = img.rotate(90, expand=True)

                    # --- Resize Image ---
                    # Resize to fit within 1920x1080 (Full HD) while maintaining aspect ratio.
                    img.thumbnail((1920, 1080))

                    # --- Standardize Format (to JPEG) ---
                    # Convert RGBA (from PNGs) to RGB for JPEG compatibility.
                    if img.mode == 'RGBA':
                        img = img.convert('RGB')
                    
                    # Save the processed image to the temporary folder.
                    new_filename = os.path.splitext(filename)[0] + ".jpg"
                    processed_filepath = os.path.join(processed_folder, new_filename)
                    img.save(processed_filepath, 'jpeg', quality=90)
                    print(f"  ✅ Processed: {filename}")

            except Exception as e:
                print(f"  ⚠️ Skipping '{filename}': Not a valid or supported image file. Error: {e}")

        # --- 3. Create PowerPoint Presentation ---
        print("\nAll images processed. Creating PowerPoint presentation...")
        
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        blank_slide_layout = prs.slide_layouts[6]

        processed_image_files = sorted(os.listdir(processed_folder))

        for image_name in processed_image_files:
            slide = prs.slides.add_slide(blank_slide_layout)
            img_path = os.path.join(processed_folder, image_name)
            
            # Add the picture and center it
            pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0))
            
            # Calculate aspect ratios
            slide_aspect_ratio = prs.slide_width / prs.slide_height
            image_aspect_ratio = pic.width / pic.height

            # Scale image to fit slide while maintaining aspect ratio
            if image_aspect_ratio > slide_aspect_ratio:
                # Image is wider than slide, scale by width
                new_width = prs.slide_width
                new_height = int(new_width / image_aspect_ratio)
            else:
                # Image is taller than slide, scale by height
                new_height = prs.slide_height
                new_width = int(new_height * image_aspect_ratio)

            # Center the image
            pic.width = new_width
            pic.height = new_height
            pic.left = int((prs.slide_width - new_width) / 2)
            pic.top = int((prs.slide_height - new_height) / 2)

        prs.save(output_path)
        print(f"\n✨ Success! Slideshow saved to your Desktop as '{output_filename}'")

    finally:
        # --- 4. Clean Up ---
        # Ensure the temporary folder is always removed.
        if os.path.exists(processed_folder):
            shutil.rmtree(processed_folder)
            print("\n🧹 Cleaned up temporary files.")

if __name__ == "__main__":
    # Ensure you have the required libraries installed:
    # pip install python-pptx Pillow
    create_slideshow_from_folder()

